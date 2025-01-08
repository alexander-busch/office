# -*- coding: utf-8 -*-
"""
Read the content from a set of pptx files and write it into a single word file.

06.03.2023
@author: alexander.busch@alumni.ntnu.no
"""

import os
import xlsxwriter
import openpyxl
import pptx
from docx import Document
from docx.shared import Pt
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.text import WD_COLOR_INDEX
from docx.shared import RGBColor
from pathlib import Path


# -----------------------------------------------------------------------------
# Parameters
# -----------------------------------------------------------------------------

# Root path to files
pathToFiles=r'C:\analysis'

# Variants
# 6 = aaa
# 8 = bbb
# 20 = ccc
variant='20'

# Local directories for pdf and pptx
if variant == '6':
    path_pptx = r'C:\aaa\Slides'    
    path_pdf = r'C:\aaa\PDF-slides'
elif variant == '8':
    path_pptx = r'C:\bbb'
    path_pdf = path_pptx
elif variant == '20':
    path_pptx = r'C:\ddd'
    path_pdf = path_pptx
else:
    print('Incorrect variant')


# -----------------------------------------------------------------------------
# Function definition
# Check if characters are valid for xml 
# https://stackoverflow.com/questions/8733233/filtering-out-certain-bytes-in-python
# -----------------------------------------------------------------------------
def valid_xml_char_ordinal(c):
    codepoint = ord(c)
    # conditions ordered by presumed frequency
    return (
        0x20 <= codepoint <= 0xD7FF or
        codepoint in (0x9, 0xA, 0xD) or
        0xE000 <= codepoint <= 0xFFFD or
        0x10000 <= codepoint <= 0x10FFFF
        )

# -----------------------------------------------------------------------------
# Read the presentation names based on the published pdf files
# -----------------------------------------------------------------------------

path = Path(os.path.join(pathToFiles, 'pdf_titles.xlsx'))

if path.is_file() is not True:
    
    # Get all pdf file names in the folder
    pdf_file_names = [file_name for file_name in os.listdir(path_pdf) if file_name.endswith('.pdf')]
    
    # Set the output file name
    output_file_name = "pdf_titles.xlsx"
    
    # Create a new workbook
    workbook = xlsxwriter.Workbook(os.path.join(pathToFiles, output_file_name))
    
    # Add a new worksheet
    worksheet = workbook.add_worksheet()
    
    # Write the file names into the worksheet
    for i, file_name in enumerate(pdf_file_names):
        worksheet.write(i, 0, file_name)
    
    # Save the workbook
    workbook.close()


# -----------------------------------------------------------------------------
# Extract all headings and textual content from the pptx presentations
# -----------------------------------------------------------------------------

# Load the Excel workbook
workbook = openpyxl.load_workbook(os.path.join(pathToFiles,'Uebersicht.xlsx'))

# Select the worksheet to read from
worksheet = workbook[variant]

# Define the range of cells to read from
text=['D3:D', str(worksheet.max_row)]
pptx_presentationname_range = ''.join(text)
text=['E3:E', str(worksheet.max_row)]
pptx_filename_range = ''.join(text)

# Use a list comprehension to read the values in the cell range
pptx_filenames = [cell[0].value for cell in worksheet[pptx_filename_range]]
pptx_presentationnames = [cell[0].value for cell in worksheet[pptx_presentationname_range]]

# Print the cell values
#print(pptxpresentations)
#print(pptxfiles)

# Create a new Word document
file_name='from_pptx.docx'
file_path = Path(os.path.join(pathToFiles, file_name))
if file_path.is_file():
    os.remove(file_path)
doc = Document()

# Try
# doc_template = os.path.join(pathToFiles, 'template.docx')
# doc = Document(doc_template)

# Retrieve and modify the Heading 1 style
style_heading1=doc.styles['Heading 1']
style_heading1.font.size = Pt(20)
style_heading1.font.name = 'Arial'
style_heading1.font.bold = True
style_heading1.font.italic = False
#style_heading1.font.color.rgb = WD_COLOR_INDEX.BLACK
style_heading1.font.color.rgb = RGBColor(0, 0, 0)


# Retrieve and modify the Heading 2 style
style_heading2=doc.styles['Heading 2']
style_heading2.font.size = Pt(15)
style_heading2.font.name = 'Arial'
style_heading2.font.bold = False
style_heading2.font.italic = False
#style_heading2.font.color.rgb = WD_COLOR_INDEX.BLACK
style_heading2.font.color.rgb = RGBColor(0, 0, 0)

# Create a new style for the slide text
style_Normal=doc.styles['Normal']
style_Normal.font.size = Pt(10)
style_Normal.font.name = 'Arial'
style_Normal.font.bold = False
style_Normal.font.italic = False

# Create a new style for the slide comment
style_Comment=doc.styles.add_style('Normal_Comment', WD_STYLE_TYPE.PARAGRAPH)
style_Comment.font.size = Pt(8)
style_Comment.font.name = 'Arial'
style_Comment.font.bold = False
style_Comment.font.italic = True



# Initialize index
index = 0

len(pptx_filenames)

# Loop all pptx files and extract data
for pptx_file in pptx_filenames:
       
    # Load current pptx file
    # parentfolder=os.path.dirname(os.path.dirname(path_pdf))
    ppt = pptx.Presentation(os.path.join(path_pptx,pptx_file))

    # Create presentation name
    if len(pptx_filenames) > 1:
        paragraph=doc.add_heading(pptx_presentationnames[index], level=1)
        paragraph.style = style_heading1
    else:
        text = ['Presentation name: ', pptx_presentationnames[index]]
        paragraph=doc.add_paragraph(''.join(text))
        paragraph.style = style_Normal
    
    # Add file name as text
    text = ['File name: ', pptx_file]
    paragraph=doc.add_paragraph(''.join(text))
    paragraph.style = style_Normal
    
    # Initialize index
    index2 = 0
    
    # Loop through each slide in the pptx file
    for slide in ppt.slides:
        
        # Check if slide heading exists
        if slide.shapes.title is not None:
            
            # Extract slide heading
            slide_heading = slide.shapes.title.text
            
            # Clean text for usage in xml
            slide_heading_cleaned = ''.join(c for c in slide_heading if valid_xml_char_ordinal(c))
            
        else:
            slide_heading_cleaned='Slide without title'
        
        # Add presentation name as text
        if index2==0:
            text = ['Presentation name: ', slide_heading_cleaned]
            paragraph=doc.add_paragraph(''.join(text))
            paragraph.style = style_Normal
        
        # Create Word heading with slide heading text
        if len(pptx_filenames) > 1:
            paragraph=doc.add_heading(slide_heading_cleaned, level=2)
            paragraph.style = style_heading2
        else:
            
            # check if slide is a section header slide 
            if slide.slide_layout.name == 'Abschnittsüberschrift':
                paragraph=doc.add_heading(slide_heading_cleaned, level=1)
                paragraph.style = style_heading1
            else:
                paragraph=doc.add_heading(slide_heading_cleaned, level=2)
                paragraph.style = style_heading2
        
        # Get slide content as plain text
        slide_content = ''
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_content += shape.text
        
        # Clean text for usage in xml
        slide_content_cleaned = ''.join(c for c in slide_content if valid_xml_char_ordinal(c))
        
        # Add cleaned slide content as plain text
        #paragraph=doc.add_paragraph(slide_content_cleaned)
        #paragraph.style = style_Normal
        
        # Extract the notes from the slide
        notes_slide = slide.notes_slide
        if hasattr(notes_slide, "text"):
            note_content = notes_slide.notes_text_frame.text
            
            # Clean text for usage in xml
            note_content_cleaned = ''.join(c for c in note_content if valid_xml_char_ordinal(c))
            
            # Add cleaned slide content as plain text
            #paragraph=doc.add_paragraph(note_content_cleaned)
            #paragraph.style = style_Comment
        
        # Increase index
        index2 += 1
    
    # Increase index
    index += 1

# Save Word document
doc.save(os.path.join(pathToFiles,file_name))
text=[variant, '.docx']
os.rename(os.path.join(pathToFiles,file_name), os.path.join(pathToFiles,''.join(text)))
