# -*- coding: utf-8 -*-
"""
Created on Fri Mar 31 15:31:52 2023

@author: alexander.busch@alumni.ntnu.no
"""

import docx
from pptx import Presentation


# Path to the master .potx file
master_file = r"C:\aaa\template.pptx"

# Function to check the indices of placeholders for a specific layout, use for debugging only
def checkPlaceholderArrayIndex(master_file, layout):
    # Create a new Presentation object using the master file
    pres = Presentation(master_file)

    # Add a new slide to the presentation using a layout by name
    slide = pres.slides.add_slide(layout)
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))

# Open Word document
doc = docx.Document(r"C:\aaa\somefile.docx")

# Path to the master .potx file
master_file = r"C:\bbb\template.potx"
master_file = r"C:\bbb\template.pptx"


# Create PowerPoint presentation
pres = Presentation(master_file)

# Get slide layouts by name
title_layout = pres.slide_layouts.get_by_name("Title")
section_layout = pres.slide_layouts.get_by_name("Section Header")
content_layout = pres.slide_layouts.get_by_name("Text_or_other")

# Function to check the indices of placeholders for a specific layout, use for debugging only
checkPlaceholderArrayIndex(master_file, title_layout)
checkPlaceholderArrayIndex(master_file, section_layout)
checkPlaceholderArrayIndex(master_file, content_layout)

section_text = ""
slide = None

# Extract heading levels and paragraphs from Word document
for para in doc.paragraphs:
    if para.style.name == "Heading 1":
        
        if section_text:
            # Add section text to slide
            slide.placeholders[14].text = section_text
            
        # Create chapter slide
        slide = pres.slides.add_slide(section_layout)
        slide.shapes.title.text = para.text.strip()
        
        section_text = ""
        
    elif para.style.name == "Heading 2":
        
        if section_text:
            # Add section text to slide
            slide.placeholders[14].text = section_text
        
        # Create section slide
        slide = pres.slides.add_slide(content_layout)
        slide.shapes.title.text = para.text.strip()
        # Extract paragraphs following heading 2
        section_text = ""
        # if hasattr(para, "next_sibling"):
        #     next_para = para.next_sibling
        #     print("Paragraph has next sibling!")
        # else:
        #     next_para = None
        #     print("Paragraph does not have next sibling.")
        # while next_para is not None and next_para.style.name != "Heading 2":
        #     section_text += next_para.text.strip() + "\n"
        #     next_para = next_para.next_sibling
        # # Add section text to slide
        # slide.placeholders[14].text = section_text
    elif para.style.name == "Normal":
        section_text += para.text.strip() + "\n"
    else:
       print('tbd')
       
# Save PowerPoint presentation
pres.save(r"C:\ccc\example.pptx")


